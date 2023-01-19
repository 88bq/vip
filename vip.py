from os import system
from os import unlink
import os

try:
    import fitz
except:
    system('pip install pymupdf')
try:
    import PyPDF2
except:
    system('pip install PyPDF2')
try:
    import gtts
except:
    system('pip install gtts')
try:
    import pytesseract
except:
    system('pip install pytesseract')
# 1788351105:AAHfYx9EsreTguPB1PUkDzHZ1ZhZz5m-nuc
try:
    import PIL
except:
    system('pip install PIL')
try:
    import io
except:
    system('pip install io')
try:
    import wand
except:
    system('pip install wand')
import shutil

try:
    from pptx import Presentation
except:
    system('pip install python-pptx')
try:
    from docx import Document
except:
    system('pip install python-docx')
from docx.enum.style import WD_STYLE_TYPE
from docx.shared import Pt, RGBColor, Inches
from docx.oxml.ns import qn
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml.shared import OxmlElement
from pptx.util import Cm

try:
    from selenium import webdriver
except:
    system('pip install selenium')
# import fitz
try:
    import telebot
except:
    system('pip install pytelegrambotapi')
from deep_translator import GoogleTranslator
import telebot
from telebot import types
import requests
from PIL import Image
import io
import pytesseract
# from wand.image import Image as wi
from pptx import Presentation
from docx import Document
import glob
import requests
import shutil
from pptx.enum.shapes import PP_MEDIA_TYPE

# pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
c = []
g = []
d = []
h = []
work = []
idi = 209501902
open('tran.txt', 'a').write(f'')
r1 = requests.session()
token = '5392900812:AAET_LLAkW4jsd6HNMKhH9fLv2Ee9G0hvoQ'
bot = telebot.TeleBot(token)


@bot.message_handler(commands=['start'])
def key(msg):
    ch = msg.chat.id
    if ch in c:
        pass
    else:
        c.append(ch)
    if msg.chat.id == idi:
        admin = types.ReplyKeyboardMarkup()
        akoky = types.KeyboardButton('الاذاعة')
        zkoky = types.KeyboardButton('الاعضاء')
        admin.add(akoky, zkoky)
        bot.reply_to(msg, 'اهلا و سهلا بك ايها المطور {}'.format(msg.from_user.first_name), reply_markup=admin)

    else:
        idd = msg.chat.id
        res = \
            r1.get(f'https://api.telegram.org/bot{token}/getChatMember?chat_id=@akokybot&user_id={idd}').json()[
                'result'][
                'status']
        if res == 'left' and msg.chat.id not in [idi, 1490464385]:
            bot.send_message(ch, 'يجب عليك الاشتراك في القناة اول و من ثم استخدم البوت \nلطفا💚\n@akokybot')
        else:
            q = types.InlineKeyboardMarkup()
            q1 = types.InlineKeyboardButton('اضغط هنا \nاذا كان هنالك مصطلح طبي تريد ترجمته', callback_data='q1')
            q.add(q1)
            a9 = types.InlineKeyboardButton("المطور 💚", url="https://t.me/Q5QQQQ")
            a10 = types.InlineKeyboardButton("قناه الاعلانات↗️", url="https://t.me/akokybot")
            q.add(a9)
            q.add(a10)
            bot.send_message(msg.chat.id, 'اهلا و سهلا بك في بالبوت الترجمة 💚', reply_markup=q)
        if f'{ch}\n' in open('tran.txt', 'r'):
            pass
        else:
            if len(open('tran.txt', 'r').read().split('\n')) > 25:
                bot.send_message(ch, 'لقد وصل البوت الى عدد المطلوب من المستخدمين و هو 25')
            else:
                open('tran.txt', 'a').write(f'{ch}\n')


@bot.message_handler(content_types='text')
def an(msg):
    ch = msg.chat.id
    x = msg.text
    try:
        if msg.chat.type != "private":
            if ch in g:
                if f'{msg.chat.id}\n' in open('tran.txt', 'r'):
                    pass
                else:
                    open('tran.txt', 'a').write(f'{msg.chat.id}\n')
                try:
                    if msg.text == 'الاعضاء' and msg.chat.id == idi:
                        x = open('tran.txt', 'r').readlines()
                        bot.reply_to(msg, '{}'.format(len(x)))
                except:
                    pass
                try:
                    if msg.text == 'الاذاعة' and msg.chat.id == idi:
                        markup = types.ForceReply(selective=False)
                        bot.send_message(msg.chat.id, "ارسل اذاعتك", reply_markup=markup)
                except:
                    pass
                try:
                    if msg.reply_to_message.text == "ارسل اذاعتك":
                        try:
                            s = msg.text
                            x = open('tran.txt', 'r')
                            for i in x:
                                try:
                                    bot.send_message(i.replace('\n', ''), s)
                                except:
                                    pass
                        except:
                            pass
                except:
                    pass
                w = msg.text.lower()
                try:
                    if msg.reply_to_message.text == "ارسل المصطلح الطبي المراد ترجمته":
                        if w[0] == 'ا' or w[0] == 'أ' and w[1] == 'ل':
                            url = f'https://context.reverso.net/translation/arabic-english/{w}'
                            head = {
                                'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8,application/signed-exchange;v=b3;q=0.9',
                                'Accept-Encoding': 'gzip, deflate, br',
                                'Accept-Language': 'en-US,en;q=0.9',
                                'Connection': 'keep-alive',
                                'Cookie': 'didomi_token=eyJ1c2VyX2lkIjoiMTc1YWYxY2QtM2E4NS02NWNhLTlhNWYtNWMzMTMxYmY0OTQyIiwiY3JlYXRlZCI6IjIwMjAtMTEtMDlUMjI6MjU6MDMuODM3WiIsInVwZGF0ZWQiOiIyMDIwLTExLTA5VDIyOjI1OjAzLjgzN1oiLCJ2ZXJzaW9uIjpudWxsfQ==; __qca=P0-2069518181-1604960703577; _ga=GA1.2.1022455205.1604960704; _fbp=fb.1.1604960704785.1696161273; __gads=ID=6a415d82c09e7f2f:T=1604961014:S=ALNI_MYnjNMXo7nFRuuahoJIpqXamOQFvA; CTXTNODEID=bstweb12; experiment_context_N7gT3vKzX=0; JSESSIONID=kamXv4rTSPJG1mzG5D0fBsyY.bst-web12; context.lastpair=ar-en; history_entry=psycho]#[gastric]#[{w}; history_pair=en-ar]#[en-ar]#[ar-en; experiment_context_E3de3pqAZ=1; _gid=GA1.2.1607272187.1606403162; context.dapppromotion-count2=1; context.dapppromotion2=0; _gat_gtag_UA_2834324_41=1',
                                'Host': 'context.reverso.net',
                                'Referer': 'https://context.reverso.net/translation/arabic-english/%7Bw%7D',
                                'Sec-Fetch-Dest': 'document',
                                'Sec-Fetch-Mode': 'navigate',
                                'Sec-Fetch-Site': 'same-origin',
                                'Sec-Fetch-User': '?1',
                                'Upgrade-Insecure-Requests': '1',
                                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/86.0.4240.198 Safari/537.36'
                            }
                            j = r1.get(url=url, headers=head)
                            v = (j.text[
                                 j.text.find('<button class="other-content" data-other="0" data-negative="') + len(
                                     '<button class="other-content" data-other="0" data-negative="'):j.text.find(
                                     'Other translations</button>')])[:-2]
                            c = (v.replace("-{", "").replace("}", "\n"))
                            bot.send_message(msg.chat.id, f'* {c} *')

                        else:
                            try:
                                if f'{w[0]}' in ['a', 'b', 'c', 'd', 'e', 'f', 'g', 'h', 'i', 'j', 'k', 'l', 'm', 'n',
                                                 'o',
                                                 'p',
                                                 'q',
                                                 'r',
                                                 's', 't', 'u', 'v', 'w', 'x', 'y', 'z']:
                                    d = w
                                    v = f'https://www.tbeeb.net/%D9%82%D8%A7%D9%85%D9%88%D8%B3-%D8%B7%D8%A8%D9%8A/search.php?q={d}+&dictionary=%D8%A8%D8%AD%D8%AB'
                                    try:
                                        i = 1
                                        zz = ''
                                        while i < len(r1.get(v).text.split('"tden">')):
                                            if i > 7:
                                                pass
                                            else:
                                                x1 = r1.get(v).text.split('"tden">')[i][
                                                     :r1.get(v).text.split('"tden">')[i].find('</td>')]
                                                x2 = r1.get(v).text.split('"tdar">')[i][
                                                     :r1.get(v).text.split('"tdar">')[i].find('</td>')]
                                                zz += f'{x2}\n{x1}\n'
                                            i += 1

                                        bot.send_message(ch, zz)
                                    except:
                                        try:
                                            url = f'https://context.reverso.net/translation/arabic-english/{w}'
                                            head = {
                                                'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8,application/signed-exchange;v=b3;q=0.9',
                                                'Accept-Encoding': 'gzip, deflate, br',
                                                'Accept-Language': 'en-US,en;q=0.9',
                                                'Connection': 'keep-alive',
                                                'Cookie': 'didomi_token=eyJ1c2VyX2lkIjoiMTc1YWYxY2QtM2E4NS02NWNhLTlhNWYtNWMzMTMxYmY0OTQyIiwiY3JlYXRlZCI6IjIwMjAtMTEtMDlUMjI6MjU6MDMuODM3WiIsInVwZGF0ZWQiOiIyMDIwLTExLTA5VDIyOjI1OjAzLjgzN1oiLCJ2ZXJzaW9uIjpudWxsfQ==; __qca=P0-2069518181-1604960703577; _ga=GA1.2.1022455205.1604960704; _fbp=fb.1.1604960704785.1696161273; __gads=ID=6a415d82c09e7f2f:T=1604961014:S=ALNI_MYnjNMXo7nFRuuahoJIpqXamOQFvA; CTXTNODEID=bstweb12; experiment_context_N7gT3vKzX=0; JSESSIONID=kamXv4rTSPJG1mzG5D0fBsyY.bst-web12; context.lastpair=ar-en; history_entry=psycho]#[gastric]#[{w}; history_pair=en-ar]#[en-ar]#[ar-en; experiment_context_E3de3pqAZ=1; _gid=GA1.2.1607272187.1606403162; context.dapppromotion-count2=1; context.dapppromotion2=0; _gat_gtag_UA_2834324_41=1',
                                                'Host': 'context.reverso.net',
                                                'Referer': 'https://context.reverso.net/translation/arabic-english/%7Bw%7D',
                                                'Sec-Fetch-Dest': 'document',
                                                'Sec-Fetch-Mode': 'navigate',
                                                'Sec-Fetch-Site': 'same-origin',
                                                'Sec-Fetch-User': '?1',
                                                'Upgrade-Insecure-Requests': '1',
                                                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/86.0.4240.198 Safari/537.36'
                                            }
                                            j = r1.get(url=url, headers=head)
                                            v = (j.text[
                                                 j.text.find(
                                                     '<button class="other-content" data-other="0" data-negative="') + len(
                                                     '<button class="other-content" data-other="0" data-negative="'):j.text.find(
                                                     'Other translations</button>')])[:-2]
                                            c = (v.replace("-{", "").replace("}", "\n"))
                                            bot.send_message(msg.chat.id, f'* {c} *')
                                        except:
                                            pass
                                    head = {
                                        'accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8,application/signed-exchange;v=b3;q=0.9',
                                        'accept-encoding': 'gzip, deflate, br',
                                        'accept-language': 'en-US,en;q=0.9',
                                        'cookie': 'CGIC=IocBdGV4dC9odG1sLGFwcGxpY2F0aW9uL3hodG1sK3htbCxhcHBsaWNhdGlvbi94bWw7cT0wLjksaW1hZ2UvYXZpZixpbWFnZS93ZWJwLGltYWdlL2FwbmcsKi8qO3E9MC44LGFwcGxpY2F0aW9uL3NpZ25lZC1leGNoYW5nZTt2PWIzO3E9MC45; HSID=ACpCi--xKNWBuzc1V; SSID=AxTv5nnKkxKd-t24h; APISID=OvtsQIlYX38DaDAB/AejTrmzLNAmxJKW5L; SAPISID=VictV1NMEEknSmKC/Ap3tCyc9Rw-nNhkKN; __Secure-3PAPISID=VictV1NMEEknSmKC/Ap3tCyc9Rw-nNhkKN; CONSENT=YES+IQ.ar+201908; SEARCH_SAMESITE=CgQIoJAB; OTZ=5713657_44_44__44_; SID=3gekKKpt2BDYs3HGt2dbGMJygRXv0S_uH30dI87nwwd9kPVGeJPLKTJluqScaUfUXXkqGg.; __Secure-3PSID=3gekKKpt2BDYs3HGt2dbGMJygRXv0S_uH30dI87nwwd9kPVGofFKYNMh0FN74vl_7RXMLg.; ANID=AHWqTUl-D-kiT0ekElpQYvwaLkYpmfjEj0_vWVNyeQ9lR3vPGLYG1ea_vAtrh9jw; 1P_JAR=2020-11-26-23; NID=204=Nvh8sYQiAGl2xGl9Rivw7mwHK0oiirbBuBjelwbKer7spRkQS3Xh-9GYO5GR2WOwUG-GICaT4wz4vP8fWCpH8Zao3AEmb_XGvZ7CoLvMiTCbcSnMjY6GCGnFcwc9Ap2D0Iu1ltcyj4qzLl25BPgZf0KdrrpVYB_UyBb_LrJWA7A4dgrtBLEwtyBFzHe1XdYKS_yOA6QzeeHlNNtZm2YYpYUzvzevK2wxD2EMs9f9OOnq_es; DV=Q77RAgvgUI1KQP_f7OZE3e4dEcltYJdw4P5MyRCtgwEAAHAXk3Gqw6shngAAAAyZb21iQXUwRQAAAA; SIDCC=AJi4QfGRLCnLLJayzTS67kgRZvvsMJ8WJx4Z3VQIypirxBw8XSvY82_9ydx6j9Wr0SSSgWwkPw; __Secure-3PSIDCC=AJi4QfFg7EgRAZ_f0vxbunVqGgMst8xwoewWB6J0ZP6y8sihJ3OufFh34ZiNnfFg87jxfwhwgg',
                                        'referer': 'https://www.google.com/',
                                        'sec-fetch-dest': 'document',
                                        'sec-fetch-mode': 'navigate',
                                        'sec-fetch-site': 'same-origin',
                                        'sec-fetch-user': '?1',
                                        'upgrade-insecure-requests': '1',
                                        'user-agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/86.0.4240.198 Safari/537.36',
                                        'x-client-data': 'CIy2yQEIpLbJAQjBtskBCKmdygEIlqzKAQisx8oBCPbHygEI6cjKAQi0y8oBCI3PygEI3NXKAQjul8sBCJGZywEImJrLARiKwcoB'
                                    }
                                    r = requests.get(
                                        f'https://www.google.com/search?q={d}+anatomy&source=lnms&tbm=isch&sa=X&ved=2ahUKEwjfv9Xg5YrvAhWHnxQKHe-2AfEQ_AUoAXoECBQQAw&biw=1821&bih=876#imgrc=3PIl4efjvoyhfM',
                                        headers=head)
                                    i = 0
                                    import shutil
                                    try:
                                        while i < 3:
                                            try:
                                                x = r.text.split('jpg"')[i].split('["http')[-1]
                                                v = f'http{x}jpg'

                                                with open(fr'koky{i}.jpg', 'wb') as mp3:
                                                    shutil.copyfileobj(r1.get(v, stream=True).raw, mp3)
                                                    bot.send_photo(chat_id=ch,
                                                                   photo=open(f'koky{i}.jpg', 'rb'))

                                            except:
                                                pass
                                            i += 1
                                    except:
                                        pass


                                else:
                                    url = f'https://context.reverso.net/translation/arabic-english/{w}'
                                    head = {
                                        'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8,application/signed-exchange;v=b3;q=0.9',
                                        'Accept-Encoding': 'gzip, deflate, br',
                                        'Accept-Language': 'en-US,en;q=0.9',
                                        'Connection': 'keep-alive',
                                        'Cookie': 'didomi_token=eyJ1c2VyX2lkIjoiMTc1YWYxY2QtM2E4NS02NWNhLTlhNWYtNWMzMTMxYmY0OTQyIiwiY3JlYXRlZCI6IjIwMjAtMTEtMDlUMjI6MjU6MDMuODM3WiIsInVwZGF0ZWQiOiIyMDIwLTExLTA5VDIyOjI1OjAzLjgzN1oiLCJ2ZXJzaW9uIjpudWxsfQ==; __qca=P0-2069518181-1604960703577; _ga=GA1.2.1022455205.1604960704; _fbp=fb.1.1604960704785.1696161273; __gads=ID=6a415d82c09e7f2f:T=1604961014:S=ALNI_MYnjNMXo7nFRuuahoJIpqXamOQFvA; CTXTNODEID=bstweb12; experiment_context_N7gT3vKzX=0; JSESSIONID=kamXv4rTSPJG1mzG5D0fBsyY.bst-web12; context.lastpair=ar-en; history_entry=psycho]#[gastric]#[{w}; history_pair=en-ar]#[en-ar]#[ar-en; experiment_context_E3de3pqAZ=1; _gid=GA1.2.1607272187.1606403162; context.dapppromotion-count2=1; context.dapppromotion2=0; _gat_gtag_UA_2834324_41=1',
                                        'Host': 'context.reverso.net',
                                        'Referer': 'https://context.reverso.net/translation/arabic-english/%7Bw%7D',
                                        'Sec-Fetch-Dest': 'document',
                                        'Sec-Fetch-Mode': 'navigate',
                                        'Sec-Fetch-Site': 'same-origin',
                                        'Sec-Fetch-User': '?1',
                                        'Upgrade-Insecure-Requests': '1',
                                        'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/86.0.4240.198 Safari/537.36'
                                    }
                                    j = r1.get(url=url, headers=head)
                                    v = (j.text[
                                         j.text.find(
                                             '<button class="other-content" data-other="0" data-negative="') + len(
                                             '<button class="other-content" data-other="0" data-negative="'):j.text.find(
                                             'Other translations</button>')])[:-2]
                                    c = (v.replace("-{", "").replace("}", "\n"))
                                    bot.send_message(msg.chat.id, f'* {c} *')
                            except:
                                pass
                except:
                    pass
                if msg.reply_to_message:
                    pass
                elif 'لغاء' in msg.text:
                    if ch in g:
                        g.remove(ch)
                    else:
                        pass
                else:
                    try:
                        url = 'https://www.arabtran.com/gtranslate/'

                        head = {
                            'accept': '*/*',
                            'accept-encoding': 'gzip, deflate, br',
                            'accept-language': 'en-US,en;q=0.9',
                            'content-length': '31',
                            'content-type': 'application/x-www-form-urlencoded; charset=UTF-8',
                            'cookie': '_ga=GA1.2.189101859.1607603223; _gid=GA1.2.154684385.1607603223; __gads=ID=ba2b5a0fe52a4c1b-22c3694088a60095:T=1607603223:RT=1607603223:S=ALNI_MbdZ0H2MeiuATaH2_vAl4hrUnPm8Q; _gat=1; sc_is_visitor_unique=rx12068393.1607603329.CA2EB2FDB25A4FB19C08D5DA0B092EBC.1.1.1.1.1.1.1.1.1',
                            'origin': 'https://www.arabtran.com',
                            'referer': 'https://www.arabtran.com/tarjamat_anjilizi_earabiun/',
                            'sec-fetch-dest': 'empty',
                            'sec-fetch-mode': 'cors',
                            'sec-fetch-site': 'same-origin',
                            'user-agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/87.0.4280.88 Safari/537.36',
                            'x-requested-with': 'XMLHttpRequest',
                        }

                        q = types.InlineKeyboardMarkup()
                        q1 = types.InlineKeyboardButton('اضغط هنا في حيال وجود مصطلح طبي 💟', callback_data='q1')
                        q.add(q1)
                        if f'{w[0]}' in ['a', 'b', 'c', 'd', 'e', 'f', 'g', 'h', 'i', 'j', 'k', 'l', 'm', 'n', 'o', 'p',
                                         'q', 'r',
                                         's', 't', 'u', 'v', 'w', 'x', 'y', 'z'] or f'{w[1]}' in ['a', 'b', 'c', 'd',
                                                                                                  'e',
                                                                                                  'f', 'g',
                                                                                                  'h', 'i', 'j', 'k',
                                                                                                  'l',
                                                                                                  'm', 'n',
                                                                                                  'o', 'p', 'q', 'r',
                                                                                                  's',
                                                                                                  't', 'u',
                                                                                                  'v', 'w', 'x', 'y',
                                                                                                  'z'] or f'{w[2]}' in [
                            'a',
                            'b',
                            'c',
                            'd',
                            'e',
                            'f',
                            'g',
                            'h',
                            'i',
                            'j',
                            'k',
                            'l',
                            'm',
                            'n',
                            'o',
                            'p',
                            'q',
                            'r',
                            's',
                            't',
                            'u',
                            'v',
                            'w',
                            'x',
                            'y',
                            'z']:
                            data = {
                                'text': w,
                                'gfrom': 'en',
                                'gto': 'ar',
                                'key': 'ABC'
                            }
                            x = r1.post(url=url, data=data, ).text
                            bot.send_message(msg.chat.id, x, reply_markup=q)

                        else:
                            data = {
                                'text': w,
                                'gfrom': 'ar',
                                'gto': 'en',
                                'key': 'ABC'
                            }
                            x = r1.post(url=url, data=data, ).text
                            bot.send_message(msg.chat.id, x, reply_markup=q)


                    except:
                        pass
            else:
                if 'تفعيل' in msg.text:
                    bot.reply_to(msg, 'تم تفعيل بوت الترجمة بنجاح 💚')
                    g.append(ch)
        else:
            if f'{msg.chat.id}\n' in open('tran.txt', 'r'):
                pass
            else:
                open('tran.txt', 'a').write(f'{msg.chat.id}\n')
            idd = msg.chat.id
            res = \
                r1.get(f'https://api.telegram.org/bot{token}/getChatMember?chat_id=@akokybot&user_id={idd}').json()[
                    'result'][
                    'status']
            if res == 'left' and msg.chat.id not in [idi, 1490464385]:
                bot.send_message(ch, 'يجب عليك الاشتراك في القناة اول و من ثم استخدم البوت \nلطفا💚\n@akokybot')
            else:
                try:
                    if msg.text == 'الاعضاء' and msg.chat.id in [idi, 1490464385]:
                        x = open('tran.txt', 'r').readlines()
                        bot.reply_to(msg, '{}'.format(len(x)))
                except:
                    pass
                try:
                    if msg.text == 'الاذاعة' and msg.chat.id == idi:
                        markup = types.ForceReply(selective=False)
                        bot.send_message(msg.chat.id, "ارسل اذاعتك", reply_markup=markup)
                except:
                    pass
                try:
                    if msg.reply_to_message.text == "ارسل اذاعتك":
                        try:
                            s = msg.text
                            x = open('tran.txt', 'r')
                            for i in x:
                                try:
                                    bot.send_message(i.replace('\n', ''), s)
                                except:
                                    pass
                        except:
                            pass
                except:
                    pass
                w = msg.text.lower()
                try:
                    if msg.reply_to_message.text == "ارسل المصطلح الطبي المراد ترجمته":
                        if w[0] == 'ا' or w[0] == 'أ' and w[1] == 'ل':
                            url = f'https://context.reverso.net/translation/arabic-english/{w}'
                            head = {
                                'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8,application/signed-exchange;v=b3;q=0.9',
                                'Accept-Encoding': 'gzip, deflate, br',
                                'Accept-Language': 'en-US,en;q=0.9',
                                'Connection': 'keep-alive',
                                'Cookie': 'didomi_token=eyJ1c2VyX2lkIjoiMTc1YWYxY2QtM2E4NS02NWNhLTlhNWYtNWMzMTMxYmY0OTQyIiwiY3JlYXRlZCI6IjIwMjAtMTEtMDlUMjI6MjU6MDMuODM3WiIsInVwZGF0ZWQiOiIyMDIwLTExLTA5VDIyOjI1OjAzLjgzN1oiLCJ2ZXJzaW9uIjpudWxsfQ==; __qca=P0-2069518181-1604960703577; _ga=GA1.2.1022455205.1604960704; _fbp=fb.1.1604960704785.1696161273; __gads=ID=6a415d82c09e7f2f:T=1604961014:S=ALNI_MYnjNMXo7nFRuuahoJIpqXamOQFvA; CTXTNODEID=bstweb12; experiment_context_N7gT3vKzX=0; JSESSIONID=kamXv4rTSPJG1mzG5D0fBsyY.bst-web12; context.lastpair=ar-en; history_entry=psycho]#[gastric]#[{w}; history_pair=en-ar]#[en-ar]#[ar-en; experiment_context_E3de3pqAZ=1; _gid=GA1.2.1607272187.1606403162; context.dapppromotion-count2=1; context.dapppromotion2=0; _gat_gtag_UA_2834324_41=1',
                                'Host': 'context.reverso.net',
                                'Referer': 'https://context.reverso.net/translation/arabic-english/%7Bw%7D',
                                'Sec-Fetch-Dest': 'document',
                                'Sec-Fetch-Mode': 'navigate',
                                'Sec-Fetch-Site': 'same-origin',
                                'Sec-Fetch-User': '?1',
                                'Upgrade-Insecure-Requests': '1',
                                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/86.0.4240.198 Safari/537.36'
                            }
                            j = r1.get(url=url, headers=head)
                            v = (j.text[
                                 j.text.find('<button class="other-content" data-other="0" data-negative="') + len(
                                     '<button class="other-content" data-other="0" data-negative="'):j.text.find(
                                     'Other translations</button>')])[:-2]
                            c = (v.replace("-{", "").replace("}", "\n"))
                            bot.send_message(msg.chat.id, f'* {c} *')

                        else:
                            try:
                                if f'{w[0]}' in ['a', 'b', 'c', 'd', 'e', 'f', 'g', 'h', 'i', 'j', 'k', 'l', 'm', 'n',
                                                 'o',
                                                 'p',
                                                 'q',
                                                 'r',
                                                 's', 't', 'u', 'v', 'w', 'x', 'y', 'z']:
                                    d = w
                                    v = f'https://www.tbeeb.net/%D9%82%D8%A7%D9%85%D9%88%D8%B3-%D8%B7%D8%A8%D9%8A/search.php?q={d}+&dictionary=%D8%A8%D8%AD%D8%AB'
                                    try:
                                        i = 1
                                        zz = ''
                                        while i < len(r1.get(v).text.split('"tden">')):
                                            if i > 7:
                                                pass
                                            else:
                                                x1 = r1.get(v).text.split('"tden">')[i][
                                                     :r1.get(v).text.split('"tden">')[i].find('</td>')]
                                                x2 = r1.get(v).text.split('"tdar">')[i][
                                                     :r1.get(v).text.split('"tdar">')[i].find('</td>')]
                                                zz += f'{x2}\n{x1}\n'
                                            i += 1

                                        bot.send_message(ch, zz)
                                    except:
                                        try:
                                            url = f'https://context.reverso.net/translation/arabic-english/{w}'
                                            head = {
                                                'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8,application/signed-exchange;v=b3;q=0.9',
                                                'Accept-Encoding': 'gzip, deflate, br',
                                                'Accept-Language': 'en-US,en;q=0.9',
                                                'Connection': 'keep-alive',
                                                'Cookie': 'didomi_token=eyJ1c2VyX2lkIjoiMTc1YWYxY2QtM2E4NS02NWNhLTlhNWYtNWMzMTMxYmY0OTQyIiwiY3JlYXRlZCI6IjIwMjAtMTEtMDlUMjI6MjU6MDMuODM3WiIsInVwZGF0ZWQiOiIyMDIwLTExLTA5VDIyOjI1OjAzLjgzN1oiLCJ2ZXJzaW9uIjpudWxsfQ==; __qca=P0-2069518181-1604960703577; _ga=GA1.2.1022455205.1604960704; _fbp=fb.1.1604960704785.1696161273; __gads=ID=6a415d82c09e7f2f:T=1604961014:S=ALNI_MYnjNMXo7nFRuuahoJIpqXamOQFvA; CTXTNODEID=bstweb12; experiment_context_N7gT3vKzX=0; JSESSIONID=kamXv4rTSPJG1mzG5D0fBsyY.bst-web12; context.lastpair=ar-en; history_entry=psycho]#[gastric]#[{w}; history_pair=en-ar]#[en-ar]#[ar-en; experiment_context_E3de3pqAZ=1; _gid=GA1.2.1607272187.1606403162; context.dapppromotion-count2=1; context.dapppromotion2=0; _gat_gtag_UA_2834324_41=1',
                                                'Host': 'context.reverso.net',
                                                'Referer': 'https://context.reverso.net/translation/arabic-english/%7Bw%7D',
                                                'Sec-Fetch-Dest': 'document',
                                                'Sec-Fetch-Mode': 'navigate',
                                                'Sec-Fetch-Site': 'same-origin',
                                                'Sec-Fetch-User': '?1',
                                                'Upgrade-Insecure-Requests': '1',
                                                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/86.0.4240.198 Safari/537.36'
                                            }
                                            j = r1.get(url=url, headers=head)
                                            v = (j.text[
                                                 j.text.find(
                                                     '<button class="other-content" data-other="0" data-negative="') + len(
                                                     '<button class="other-content" data-other="0" data-negative="'):j.text.find(
                                                     'Other translations</button>')])[:-2]
                                            c = (v.replace("-{", "").replace("}", "\n"))
                                            bot.send_message(msg.chat.id, f'* {c} *')
                                        except:
                                            pass
                                    head = {
                                        'accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8,application/signed-exchange;v=b3;q=0.9',
                                        'accept-encoding': 'gzip, deflate, br',
                                        'accept-language': 'en-US,en;q=0.9',
                                        'cookie': 'CGIC=IocBdGV4dC9odG1sLGFwcGxpY2F0aW9uL3hodG1sK3htbCxhcHBsaWNhdGlvbi94bWw7cT0wLjksaW1hZ2UvYXZpZixpbWFnZS93ZWJwLGltYWdlL2FwbmcsKi8qO3E9MC44LGFwcGxpY2F0aW9uL3NpZ25lZC1leGNoYW5nZTt2PWIzO3E9MC45; HSID=ACpCi--xKNWBuzc1V; SSID=AxTv5nnKkxKd-t24h; APISID=OvtsQIlYX38DaDAB/AejTrmzLNAmxJKW5L; SAPISID=VictV1NMEEknSmKC/Ap3tCyc9Rw-nNhkKN; __Secure-3PAPISID=VictV1NMEEknSmKC/Ap3tCyc9Rw-nNhkKN; CONSENT=YES+IQ.ar+201908; SEARCH_SAMESITE=CgQIoJAB; OTZ=5713657_44_44__44_; SID=3gekKKpt2BDYs3HGt2dbGMJygRXv0S_uH30dI87nwwd9kPVGeJPLKTJluqScaUfUXXkqGg.; __Secure-3PSID=3gekKKpt2BDYs3HGt2dbGMJygRXv0S_uH30dI87nwwd9kPVGofFKYNMh0FN74vl_7RXMLg.; ANID=AHWqTUl-D-kiT0ekElpQYvwaLkYpmfjEj0_vWVNyeQ9lR3vPGLYG1ea_vAtrh9jw; 1P_JAR=2020-11-26-23; NID=204=Nvh8sYQiAGl2xGl9Rivw7mwHK0oiirbBuBjelwbKer7spRkQS3Xh-9GYO5GR2WOwUG-GICaT4wz4vP8fWCpH8Zao3AEmb_XGvZ7CoLvMiTCbcSnMjY6GCGnFcwc9Ap2D0Iu1ltcyj4qzLl25BPgZf0KdrrpVYB_UyBb_LrJWA7A4dgrtBLEwtyBFzHe1XdYKS_yOA6QzeeHlNNtZm2YYpYUzvzevK2wxD2EMs9f9OOnq_es; DV=Q77RAgvgUI1KQP_f7OZE3e4dEcltYJdw4P5MyRCtgwEAAHAXk3Gqw6shngAAAAyZb21iQXUwRQAAAA; SIDCC=AJi4QfGRLCnLLJayzTS67kgRZvvsMJ8WJx4Z3VQIypirxBw8XSvY82_9ydx6j9Wr0SSSgWwkPw; __Secure-3PSIDCC=AJi4QfFg7EgRAZ_f0vxbunVqGgMst8xwoewWB6J0ZP6y8sihJ3OufFh34ZiNnfFg87jxfwhwgg',
                                        'referer': 'https://www.google.com/',
                                        'sec-fetch-dest': 'document',
                                        'sec-fetch-mode': 'navigate',
                                        'sec-fetch-site': 'same-origin',
                                        'sec-fetch-user': '?1',
                                        'upgrade-insecure-requests': '1',
                                        'user-agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/86.0.4240.198 Safari/537.36',
                                        'x-client-data': 'CIy2yQEIpLbJAQjBtskBCKmdygEIlqzKAQisx8oBCPbHygEI6cjKAQi0y8oBCI3PygEI3NXKAQjul8sBCJGZywEImJrLARiKwcoB'
                                    }
                                    r = requests.get(
                                        f'https://www.google.com/search?q={d}+anatomy&source=lnms&tbm=isch&sa=X&ved=2ahUKEwjfv9Xg5YrvAhWHnxQKHe-2AfEQ_AUoAXoECBQQAw&biw=1821&bih=876#imgrc=3PIl4efjvoyhfM',
                                        headers=head)
                                    i = 0
                                    import shutil
                                    try:
                                        while i < 3:
                                            try:
                                                x = r.text.split('jpg"')[i].split('["http')[-1]
                                                v = f'http{x}jpg'

                                                with open(fr'koky{i}.jpg', 'wb') as mp3:
                                                    shutil.copyfileobj(r1.get(v, stream=True).raw, mp3)
                                                    bot.send_photo(chat_id=ch,
                                                                   photo=open(f'koky{i}.jpg', 'rb'))

                                            except:
                                                pass
                                            i += 1
                                    except:
                                        pass


                                else:
                                    url = f'https://context.reverso.net/translation/arabic-english/{w}'
                                    head = {
                                        'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8,application/signed-exchange;v=b3;q=0.9',
                                        'Accept-Encoding': 'gzip, deflate, br',
                                        'Accept-Language': 'en-US,en;q=0.9',
                                        'Connection': 'keep-alive',
                                        'Cookie': 'didomi_token=eyJ1c2VyX2lkIjoiMTc1YWYxY2QtM2E4NS02NWNhLTlhNWYtNWMzMTMxYmY0OTQyIiwiY3JlYXRlZCI6IjIwMjAtMTEtMDlUMjI6MjU6MDMuODM3WiIsInVwZGF0ZWQiOiIyMDIwLTExLTA5VDIyOjI1OjAzLjgzN1oiLCJ2ZXJzaW9uIjpudWxsfQ==; __qca=P0-2069518181-1604960703577; _ga=GA1.2.1022455205.1604960704; _fbp=fb.1.1604960704785.1696161273; __gads=ID=6a415d82c09e7f2f:T=1604961014:S=ALNI_MYnjNMXo7nFRuuahoJIpqXamOQFvA; CTXTNODEID=bstweb12; experiment_context_N7gT3vKzX=0; JSESSIONID=kamXv4rTSPJG1mzG5D0fBsyY.bst-web12; context.lastpair=ar-en; history_entry=psycho]#[gastric]#[{w}; history_pair=en-ar]#[en-ar]#[ar-en; experiment_context_E3de3pqAZ=1; _gid=GA1.2.1607272187.1606403162; context.dapppromotion-count2=1; context.dapppromotion2=0; _gat_gtag_UA_2834324_41=1',
                                        'Host': 'context.reverso.net',
                                        'Referer': 'https://context.reverso.net/translation/arabic-english/%7Bw%7D',
                                        'Sec-Fetch-Dest': 'document',
                                        'Sec-Fetch-Mode': 'navigate',
                                        'Sec-Fetch-Site': 'same-origin',
                                        'Sec-Fetch-User': '?1',
                                        'Upgrade-Insecure-Requests': '1',
                                        'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/86.0.4240.198 Safari/537.36'
                                    }
                                    j = r1.get(url=url, headers=head)
                                    v = (j.text[
                                         j.text.find(
                                             '<button class="other-content" data-other="0" data-negative="') + len(
                                             '<button class="other-content" data-other="0" data-negative="'):j.text.find(
                                             'Other translations</button>')])[:-2]
                                    c = (v.replace("-{", "").replace("}", "\n"))
                                    bot.send_message(msg.chat.id, f'* {c} *')
                            except:
                                pass


                except:
                    pass
                if msg.reply_to_message:
                    pass
                else:
                    try:
                        url = 'https://www.arabtran.com/gtranslate/'

                        head = {
                            'accept': '*/*',
                            'accept-encoding': 'gzip, deflate, br',
                            'accept-language': 'en-US,en;q=0.9',
                            'content-length': '31',
                            'content-type': 'application/x-www-form-urlencoded; charset=UTF-8',
                            'cookie': '_ga=GA1.2.189101859.1607603223; _gid=GA1.2.154684385.1607603223; __gads=ID=ba2b5a0fe52a4c1b-22c3694088a60095:T=1607603223:RT=1607603223:S=ALNI_MbdZ0H2MeiuATaH2_vAl4hrUnPm8Q; _gat=1; sc_is_visitor_unique=rx12068393.1607603329.CA2EB2FDB25A4FB19C08D5DA0B092EBC.1.1.1.1.1.1.1.1.1',
                            'origin': 'https://www.arabtran.com',
                            'referer': 'https://www.arabtran.com/tarjamat_anjilizi_earabiun/',
                            'sec-fetch-dest': 'empty',
                            'sec-fetch-mode': 'cors',
                            'sec-fetch-site': 'same-origin',
                            'user-agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/87.0.4280.88 Safari/537.36',
                            'x-requested-with': 'XMLHttpRequest',
                        }

                        q = types.InlineKeyboardMarkup()
                        q1 = types.InlineKeyboardButton('اضغط هنا في حيال وجود مصطلح طبي 💟', callback_data='q1')
                        q.add(q1)
                        if f'{w[0]}' in ['a', 'b', 'c', 'd', 'e', 'f', 'g', 'h', 'i', 'j', 'k', 'l', 'm', 'n', 'o', 'p',
                                         'q',
                                         'r',
                                         's', 't', 'u', 'v', 'w', 'x', 'y', 'z'] or f'{w[1]}' in ['a', 'b', 'c', 'd',
                                                                                                  'e',
                                                                                                  'f',
                                                                                                  'g',
                                                                                                  'h', 'i', 'j', 'k',
                                                                                                  'l',
                                                                                                  'm',
                                                                                                  'n',
                                                                                                  'o', 'p', 'q', 'r',
                                                                                                  's',
                                                                                                  't',
                                                                                                  'u',
                                                                                                  'v', 'w', 'x', 'y',
                                                                                                  'z'] or f'{w[2]}' in [
                            'a',
                            'b',
                            'c',
                            'd',
                            'e',
                            'f',
                            'g',
                            'h',
                            'i',
                            'j',
                            'k',
                            'l',
                            'm',
                            'n',
                            'o',
                            'p',
                            'q',
                            'r',
                            's',
                            't',
                            'u',
                            'v',
                            'w',
                            'x',
                            'y',
                            'z']:
                            data = {
                                'text': w,
                                'gfrom': 'en',
                                'gto': 'ar',
                                'key': 'ABC'
                            }
                            x = GoogleTranslator(source='auto', target='ar').translate(
                                text=msg.text)
                            bot.send_message(msg.chat.id, x, reply_markup=q)

                        else:
                            data = {
                                'text': w,
                                'gfrom': 'ar',
                                'gto': 'en',
                                'key': 'ABC'
                            }
                            x = GoogleTranslator(source='auto', target='en').translate(
                                text=msg.text)

                            bot.send_message(msg.chat.id, x, reply_markup=q)

                    except:
                        pass
    except:
        pass


@bot.callback_query_handler(lambda call: True)
def any(call):
    markup = types.ForceReply(selective=False)
    ch = call.message.chat.id
    if f'{ch}\n' in open('tran.txt', 'r'):
        pass
    else:
        if len(open('tran.txt', 'r').read().split('\n')) > 25:
            bot.send_message(ch, 'لقد وصل البوت الى عدد المطلوب من المستخدمين و هو 25')
        else:
            open('tran.txt', 'a').write(f'{ch}\n')
    if f'{ch}\n' in open('tran.txt', 'r'):
        try:
            if call.data == 'q1':
                markup = types.ForceReply(selective=False)
                bot.send_message(call.message.chat.id, "ارسل المصطلح الطبي المراد ترجمته", reply_markup=markup)
            if call.data == 'n1':
                ph = call.message.photo[2].file_id
                file_info = bot.get_file(ph)
                downloaded_file = bot.download_file(file_info.file_path)

                with open('xc1.jpg', 'wb') as new_file:
                    new_file.write(downloaded_file)
                im = Image.open('xc1.jpg')
                text = pytesseract.image_to_string(im)
                # open(f'sound/{ch}.txt', 'w', encoding='utf-8').write(text)

                # txt = open(f'sound/{ch}.txt', 'r', encoding='utf-8').read()
                gtts.gTTS(text, lang='en', slow=True).save('koky.mp3')
                bot.send_audio(ch, open('koky.mp3', 'rb'))
            if call.data == 'n2':
                ph = call.message.photo[2].file_id
                file_info = bot.get_file(ph)
                downloaded_file = bot.download_file(file_info.file_path)

                with open('xc1.jpg', 'wb') as new_file:
                    new_file.write(downloaded_file)
                im = Image.open('xc1.jpg')
                text = pytesseract.image_to_string(im)
                # open(f'sound/{ch}.txt', 'w', encoding='utf-8').write(text)
                # txt = open(f'sound/{ch}.txt', 'r', encoding='utf-8').read()
                bot.send_message(ch, text)
            if call.data == 'n3':
                ph = call.message.photo[2].file_id
                file_info = bot.get_file(ph)
                downloaded_file = bot.download_file(file_info.file_path)

                with open('xc1.jpg', 'wb') as new_file:
                    new_file.write(downloaded_file)
                    im = Image.open('xc1.jpg')
                    text = pytesseract.image_to_string(im)
                    url = 'https://www.arabtran.com/gtranslate/'
                    head = {
                        'accept': '*/*',
                        'accept-encoding': 'gzip, deflate, br',
                        'accept-language': 'en-US,en;q=0.9',
                        'content-length': '31',
                        'content-type': 'application/x-www-form-urlencoded; charset=UTF-8',
                        'cookie': '_ga=GA1.2.189101859.1607603223; _gid=GA1.2.154684385.1607603223; __gads=ID=ba2b5a0fe52a4c1b-22c3694088a60095:T=1607603223:RT=1607603223:S=ALNI_MbdZ0H2MeiuATaH2_vAl4hrUnPm8Q; _gat=1; sc_is_visitor_unique=rx12068393.1607603329.CA2EB2FDB25A4FB19C08D5DA0B092EBC.1.1.1.1.1.1.1.1.1',
                        'origin': 'https://www.arabtran.com',
                        'referer': 'https://www.arabtran.com/tarjamat_anjilizi_earabiun/',
                        'sec-fetch-dest': 'empty',
                        'sec-fetch-mode': 'cors',
                        'sec-fetch-site': 'same-origin',
                        'user-agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/87.0.4280.88 Safari/537.36',
                        'x-requested-with': 'XMLHttpRequest',
                    }
                    data = {
                        'text': text,
                        'gfrom': 'en',
                        'gto': 'ar',
                        'key': 'ABC'
                    }
                    x = GoogleTranslator(source='auto', target='ar').translate(
                        text=text)
                    bot.send_message(ch, x)
            if call.data == 'p1':
                file_info = bot.get_file(call.message.document.file_id)
                downloaded_file = bot.download_file(file_info.file_path)
                try:
                    x = call.message.document.file_name[-4:]
                    if x == '.pdf' or x == '.PDF':
                        if ch in h:
                            bot.send_message(ch, 'انتضر قليلا من فضلك لكي يتم التحميل 💚')
                        else:
                            h.append(ch)
                            bot.send_message(ch, 'حسنا, سيتم التحميل الرجاء انتظار دقيقة واحدة 💚')
                            with open('koky.pdf', 'wb') as new_file:
                                new_file.write(downloaded_file)

                            with fitz.open("koky.pdf") as doc:
                                text = ""
                                for page in doc:
                                    text += page.getText()
                                    text += f'----{str(page)[:7]}-------\n'
                                    text += '|#|'
                                    # open(f'sound/{ch}.txt', 'a').write(f'{text}\n')
                                # open(f'sound/{ch}.txt', 'a', encoding='utf-8').write(text)

                                # txt = open(f'sound/{ch}.txt', 'r', encoding='utf-8').read()
                                url = 'http://www.fromtexttospeech.com/'
                                data = {
                                    'input_text': text,
                                    'language': 'US English',
                                    'voice': 'IVONA Kimberly22',
                                    'speed': '-2',
                                    'action': 'process_text'
                                }
                                a = r1.post(url=url, data=data, ).text
                                image_location_response = requests.get(
                                    f'http://www.fromtexttospeech.com/{a[a.find("output/"):][:a[a.find("output/"):].find("mp3")]}mp3',
                                    stream=True)

                                with open(fr'koky.mp3', 'wb') as mp3:
                                    shutil.copyfileobj(image_location_response.raw, mp3)
                                bot.send_audio(ch, open('koky.mp3', 'rb'))
                    else:
                        pp = []

                        x = (call.message.document.file_name).split('.')[-1]

                        if x == 'pptx' or x == 'PPT':
                            with open('Ahmed.pptx', 'wb') as new_file:
                                new_file.write(downloaded_file)

                            for eachfile in glob.glob("Ahmed.pptx"):

                                prs = Presentation(eachfile)

                                u = ''
                                for slide in prs.slides:

                                    for shape in slide.shapes:
                                        if hasattr(shape, "text"):
                                            u += shape.text

                                url = 'http://www.fromtexttospeech.com/'
                                data = {
                                    'input_text': u,
                                    'language': 'US English',
                                    'voice': 'IVONA Kimberly22',
                                    'speed': '-2',
                                    'action': 'process_text'
                                }
                                a = r1.post(url=url, data=data, ).text
                                image_location_response = requests.get(
                                    f'http://www.fromtexttospeech.com/{a[a.find("output/"):][:a[a.find("output/"):].find("mp3")]}mp3',
                                    stream=True)

                                with open(fr'koky.mp3', 'wb') as mp3:
                                    shutil.copyfileobj(image_location_response.raw, mp3)

                                bot.send_audio(ch, open('koky.mp3', 'rb'))

                        elif x == 'docx' or x == 'DOCx':
                            u = ''
                            with open('Ahmed.docx', 'wb') as new_file:
                                new_file.write(downloaded_file)
                            document = Document('Ahmed.docx')

                            for p in document.paragraphs:
                                u += p.text
                                u += '\n'
                            url = 'http://www.fromtexttospeech.com/'
                            data = {
                                'input_text': u,
                                'language': 'US English',
                                'voice': 'IVONA Kimberly22',
                                'speed': '-2',
                                'action': 'process_text'
                            }
                            a = r1.post(url=url, data=data, ).text
                            image_location_response = requests.get(
                                f'http://www.fromtexttospeech.com/{a[a.find("output/"):][:a[a.find("output/"):].find("mp3")]}mp3',
                                stream=True)

                            with open(fr'koky.mp3', 'wb') as mp3:
                                shutil.copyfileobj(image_location_response.raw, mp3)

                            bot.send_audio(ch, open('koky.mp3', 'rb'))

                        else:
                            bot.send_message(ch, 'عذرا لم يتم ترجمة ملفك لانه ليس✳️ \npdf or pptx or docx\n')

                    try:
                        h.remove(ch)
                    except:
                        pass

                except:
                    pass
            if call.data == 'p3':
                file_info = bot.get_file(call.message.document.file_id)
                downloaded_file = bot.download_file(file_info.file_path)
                x = call.message.document.file_name[-4:]
                if x == '.pdf' or x == '.PDF':

                    if ch in h:
                        bot.send_message(ch, 'انتضر قليلا من فضلك لكي يتم التحميل 💚')
                    else:
                        h.append(ch)
                        bot.send_message(ch, 'حسنا, سيتم التحميل الرجاء انتظار دقيقة واحدة 💚')
                        with open('ahmed.pdf', 'wb') as new_file:
                            new_file.write(downloaded_file)
                        doc = fitz.open(f'ahmed.pdf')
                        v = 0
                        document = Document()

                        with fitz.open("ahmed.pdf") as doc:
                            text = ""

                            v = 0

                            for page in doc:
                                text = page.getText()
                                text += f'----------------{str(page)[:7]}-------------------\n'
                                text += '|#|'
                                pix = doc.loadPage(v).getPixmap()  # number of page
                                output = f"{ch}.png"

                                pix.writePNG(output)
                                im = Image.open(output)
                                text = pytesseract.image_to_string(im)
                                tran = GoogleTranslator(source='en', target='ar').translate(
                                    text=text)
                                # bot.send_photo(ch, open(f'{ch}.png', 'rb'))
                                v += 1
                                rr = 0
                                xc = ''
                                try:
                                    a1 = text.split('\n')
                                    a2 = tran.split('\n')

                                    while rr < len(text.split('\n')):
                                        if a1[rr] == '' or a1[rr] == ' ':
                                            pass
                                        else:
                                            xc = f'{a2[rr]}\n{a1[rr]}'

                                            document.add_paragraph(f'{xc}')
                                        rr += 1
                                except:
                                    pass
                                z = f'----{str(page)[:7]}-------\n'
                                document.add_picture(f'{ch}.png')
                                document.add_paragraph(f'{z}\n')
                                # document.add_picture(f'{ch}.png')
                                # document.add_paragraph(f'{tran}')
                            document.save('Ahmed2.docx')
                            bot.send_document(ch, open(f'Ahmed2.docx', 'rb'),
                                              caption=f'{call.message.document.file_name}\nترجمة صور داخل الملفات ')
                else:
                    bot.answer_callback_query(call.id, f'غير متوفره الخاصية الان لملف من نوع {x} ⚛️', True)

                try:
                    h.remove(ch)
                except:
                    pass
            if call.data == 'p2':
                file_info = bot.get_file(call.message.document.file_id)
                downloaded_file = bot.download_file(file_info.file_path)
                try:
                    x = call.message.document.file_name[-4:]
                    if x == '.pdf' or x == '.PDF':
                        if ch in h:
                            bot.send_message(ch, 'انتضر قليلا من فضلك لكي يتم التحميل 💚')
                        else:
                            h.append(ch)
                            bot.send_message(ch, 'حسنا, سيتم التحميل الرجاء انتظار دقيقة واحدة 💚')
                            with open('koky.pdf', 'wb') as new_file:
                                new_file.write(downloaded_file)

                            with fitz.open("koky.pdf") as doc:
                                open('ahmed.txt', 'w').write('')
                                document = Document()
                                text = ""
                                for page in doc:
                                    try:
                                        for img in doc.getPageImageList(page):
                                            xref = img[0]
                                            pix = fitz.Pixmap(doc, xref)
                                            if pix.n < 5:  # this is GRAY or RGB
                                                pix.writePNG(f"{ch}.png")
                                                document.add_picture(f'{ch}.png')

                                                v += 1
                                            else:  # CMYK: convert to RGB first
                                                pix1 = fitz.Pixmap(fitz.csRGB, pix)
                                                pix1.writePNG("p%s-%s.png" % (page, xref))
                                                pix1 = None
                                            pix = None
                                    except:
                                        pass
                                    text = page.getText()

                                    rr = 0
                                    xc = ''
                                    try:
                                        a1 = text.split('\n')

                                        while rr < len(text.split('\n')):
                                            if a1[rr] == '' or a1[rr] == ' ':
                                                pass
                                            else:
                                                if len(a1[rr]) < 22:
                                                    trans = GoogleTranslator(source='en', target='ar').translate(
                                                        text=a1[rr])
                                                    xc = f'{a1[rr]} {trans}'
                                                else:
                                                    trans = GoogleTranslator(source='en', target='ar').translate(
                                                        text=a1[rr])
                                                    xc = f'{trans}\n{a1[rr]}'

                                                document.add_paragraph(f'{xc}')

                                            rr += 1
                                    except:
                                        pass
                                    z = f'----{str(page)[:7]}-------\n'

                                    document.add_paragraph(f'{z}\n')

                            # shd = OxmlElement('w:background')
                            # shd.set(qn('w:color'), '50D1000D100D')  # black color
                            # shd.set(qn('w:themeColor'), 'text1')
                            # shd.set(qn('w:themeTint'), 'F2')

                            # document.element.insert(0,shd)
                            # shd1 = OxmlElement('w:displayBackgroundShape')
                            # document.settings.element.insert(0, shd1)

                            document.save('Ahmed.docx')
                            bot.send_document(ch, open('Ahmed.docx', 'rb'), caption=call.message.document.file_name)
                    else:
                        if ch in h:
                            bot.send_message(ch, 'انتضر قليلا من فضلك لكي يتم التحميل 💚')
                        else:
                            h.append(ch)
                            bot.send_message(ch, 'حسنا, سيتم التحميل الرجاء انتظار دقيقة واحدة 💚')
                        pp = []

                        x = (call.message.document.file_name).split('.')[-1]

                        if x == 'pptx' or x == 'PPT':
                            with open('Ahmed.pptx', 'wb') as new_file:
                                new_file.write(downloaded_file)

                            for eachfile in glob.glob("Ahmed.pptx"):

                                prs = Presentation(eachfile)
                                document = Document()

                                u = ''
                                i = 0
                                for slide in prs.slides:

                                    for shape in slide.shapes:
                                        if hasattr(shape, "text"):
                                            u += shape.text
                                    u += '---------'
                                    aa = u.split('---------')

                                    trans = GoogleTranslator(source='en', target='ar').translate(
                                        text=str(aa[i]))
                                    try:
                                        ss = aa[i].split('\n')
                                        ss1 = trans.split('\n')
                                        ss2 = 0
                                        while ss2 < len(ss):
                                            document.add_paragraph(f'{ss1[ss2]}\n{ss[ss2]}\n----------------')

                                            ss2 += 1
                                    except:
                                        pass

                                    i += 1
                                # shd = OxmlElement('w:background')
                                # shd.set(qn('w:color'), '50D1000D100D')  # black color
                                # shd.set(qn('w:themeColor'), 'text1')
                                # shd.set(qn('w:themeTint'), 'F2')

                                # document.element.insert(0, shd)
                                # shd1 = OxmlElement('w:displayBackgroundShape')
                                # document.settings.element.insert(0, shd1)

                                document.save('Ahmed.docx')
                                bot.send_document(ch, open('Ahmed.docx', 'rb'), caption=call.message.document.file_name)

                        elif x == 'docx' or x == 'DOCx':
                            u = ''
                            with open('Ahmed1.docx', 'wb') as new_file:
                                new_file.write(downloaded_file)
                            document1 = Document('Ahmed1.docx')
                            document = Document()

                            for p in document1.paragraphs:
                                u = p.text
                                u += '\n'

                                if len(u) > 5000:
                                    print('ok')
                                trans = GoogleTranslator(source='en', target='ar').translate(
                                    text=str(u))

                                document.add_paragraph(f'{trans}\n{u}\n----------------')
                            shd = OxmlElement('w:background')
                            shd.set(qn('w:color'), '50D1000D100D')  # black color
                            shd.set(qn('w:themeColor'), 'text1')
                            shd.set(qn('w:themeTint'), 'F2')

                            document.element.insert(0, shd)
                            shd1 = OxmlElement('w:displayBackgroundShape')
                            document.settings.element.insert(0, shd1)

                            document.save('Ahmed.docx')
                            bot.send_document(ch, open('Ahmed.docx', 'rb'), caption=call.message.document.file_name)




                        else:
                            bot.send_message(ch, 'عذرا لم يتم ترجمة ملفك لانه ليس✳️ \npdf or pptx or docx\n')

                    try:
                        h.remove(ch)
                    except:
                        pass

                except:
                    pass
            if call.data == 'p4':
                file_info = bot.get_file(call.message.document.file_id)
                downloaded_file = bot.download_file(file_info.file_path)
                try:
                    x = call.message.document.file_name[-4:]
                    if x == '.pdf' or x == '.PDF':
                        if ch in h:
                            bot.send_message(ch, 'انتضر قليلا من فضلك لكي يتم التحميل 💚')
                        else:
                            h.append(ch)
                            bot.send_message(ch, 'حسنا, سيتم التحميل الرجاء انتظار دقيقة واحدة 💚')
                            with open('koky.pdf', 'wb') as new_file:
                                new_file.write(downloaded_file)

                            with fitz.open("koky.pdf") as doc:
                                pptx1 = Presentation()

                                text = ""
                                for page in doc:
                                    text = page.getText()
                                    url = 'http://www.fromtexttospeech.com/'
                                    data = {
                                        'input_text': text,
                                        'language': 'US English',
                                        'voice': 'IVONA Kimberly22',
                                        'speed': '-2',
                                        'action': 'process_text'
                                    }
                                    a = r1.post(url=url, data=data, ).text
                                    image_location_response = requests.get(
                                        f'http://www.fromtexttospeech.com/{a[a.find("output/"):][:a[a.find("output/"):].find("mp3")]}mp3',
                                        stream=True)

                                    with open(fr'koky.mp3', 'wb') as mp3:
                                        shutil.copyfileobj(image_location_response.raw, mp3)
                                    if len(text) > 5000:
                                        trans = ''
                                        i = len(text) // 5000

                                        ii = 0
                                        iii = 0
                                        while ii < i + 1:
                                            tran = GoogleTranslator(source='en', target='ar').translate(
                                                text=text[iii:4999 + iii])

                                            ii += 1
                                            iii += 5000
                                        trans += tran
                                    else:
                                        trans = GoogleTranslator(source='en', target='ar').translate(
                                            text=text)
                                    slade = pptx1.slides.add_slide(pptx1.slide_layouts[1])

                                    vid = slade.shapes.add_movie(movie_file=fr'koky.mp3', left=Inches(5),
                                                                 top=Inches(0), width=Inches(1), height=Inches(1),
                                                                 mime_type='audio/mp3')
                                    vid.media_type == PP_MEDIA_TYPE.SOUND
                                    rr = 0
                                    xc = ''
                                    try:
                                        a1 = text.split('\n')
                                        a2 = trans.split('\n')

                                        while rr < len(text.split('\n')):
                                            if a1[rr] == '' or a1[rr] == ' ':
                                                pass
                                            else:
                                                xc += f'{a2[rr]}\n{a1[rr]}\n'

                                            rr += 1
                                        text_box = slade.shapes.add_textbox(Cm(1), Cm(1), Cm(5), Cm(2))
                                        text_box.text = xc


                                    except:
                                        pass
                                pptx1.save('1.pptx')

                                bot.send_document(ch, open('1.pptx', 'rb'))
                    else:
                        if ch in h:
                            bot.send_message(ch, 'انتضر قليلا من فضلك لكي يتم التحميل 💚')
                        else:
                            h.append(ch)
                            bot.send_message(ch, 'حسنا, سيتم التحميل الرجاء انتظار دقيقة واحدة 💚')
                        pp = []

                        x = (call.message.document.file_name).split('.')[-1]

                        if x == 'pptx' or x == 'PPT':
                            pptx1 = Presentation()
                            with open('Ahmed.pptx', 'wb') as new_file:
                                new_file.write(downloaded_file)

                            for eachfile in glob.glob("Ahmed.pptx"):

                                prs = Presentation(eachfile)
                                document = Document()

                                u = ''
                                i = 0
                                for slide in prs.slides:

                                    for shape in slide.shapes:
                                        if hasattr(shape, "text"):
                                            u += shape.text
                                    u += '---------'
                                    aa = u.split('---------')

                                    trans = GoogleTranslator(source='en', target='ar').translate(
                                        text=str(aa[i]))
                                    url = 'http://www.fromtexttospeech.com/'
                                    data = {
                                        'input_text': aa[i],
                                        'language': 'US English',
                                        'voice': 'IVONA Kimberly22',
                                        'speed': '-2',
                                        'action': 'process_text'
                                    }
                                    a = r1.post(url=url, data=data, ).text
                                    image_location_response = requests.get(
                                        f'http://www.fromtexttospeech.com/{a[a.find("output/"):][:a[a.find("output/"):].find("mp3")]}mp3',
                                        stream=True)

                                    with open(fr'koky.mp3', 'wb') as mp3:
                                        shutil.copyfileobj(image_location_response.raw, mp3)
                                    slade = pptx1.slides.add_slide(pptx1.slide_layouts[1])

                                    vid = slade.shapes.add_movie(movie_file=fr'koky.mp3', left=Inches(5),
                                                                 top=Inches(0), width=Inches(1), height=Inches(1),
                                                                 mime_type='audio/mp3')
                                    vid.media_type == PP_MEDIA_TYPE.SOUND
                                    try:
                                        ss = aa[i].split('\n')
                                        ss1 = trans.split('\n')
                                        ss2 = 0
                                        ss3 = ''
                                        while ss2 < len(ss):
                                            ss3 += f'{ss1[ss2]}\n{ss[ss2]}\n----------------\n'
                                            ss2 += 1

                                    except:
                                        pass
                                    shapes = slade.shapes
                                    body_shape = shapes.placeholders[1]
                                    tf = body_shape.text_frame
                                    p = tf.add_paragraph()
                                    p.text = ss3
                                    p.level = 0
                                    p.font.size = Pt(15)

                                    i += 1
                                pptx1.save('1.pptx')
                                bot.send_document(ch, open('1.pptx', 'rb'))

                        elif x == 'docx' or x == 'DOCx':
                            pptx1 = Presentation()
                            u = ''
                            with open('Ahmed1.docx', 'wb') as new_file:
                                new_file.write(downloaded_file)

                            document1 = Document('Ahmed1.docx')
                            document = Document()
                            q = 0
                            ss3 = ''
                            ss4 = ''
                            for p in document1.paragraphs:
                                u = p.text

                                if len(u) > 5000:
                                    print('ok')

                                trans = GoogleTranslator(source='en', target='ar').translate(
                                    text=str(u))
                                ss3 += f'{trans}\n{u}\n'
                                ss4 += u
                                q += 1
                                if q == 6:
                                    url = 'http://www.fromtexttospeech.com/'
                                    data = {
                                        'input_text': ss4,
                                        'language': 'US English',
                                        'voice': 'IVONA Kimberly22',
                                        'speed': '-2',
                                        'action': 'process_text'
                                    }
                                    a = r1.post(url=url, data=data, ).text
                                    image_location_response = requests.get(
                                        f'http://www.fromtexttospeech.com/{a[a.find("output/"):][:a[a.find("output/"):].find("mp3")]}mp3',
                                        stream=True)

                                    with open(fr'koky.mp3', 'wb') as mp3:
                                        shutil.copyfileobj(image_location_response.raw, mp3)
                                    slade = pptx1.slides.add_slide(pptx1.slide_layouts[1])
                                    vid = slade.shapes.add_movie(movie_file=fr'koky.mp3', left=Inches(5),
                                                                 top=Inches(0), width=Inches(1), height=Inches(1),
                                                                 mime_type='audio/mp3')
                                    vid.media_type == PP_MEDIA_TYPE.SOUND
                                    shapes = slade.shapes
                                    body_shape = shapes.placeholders[1]
                                    tf = body_shape.text_frame
                                    p = tf.add_paragraph()
                                    p.text = ss3
                                    p.level = 0
                                    p.font.size = Pt(15)
                                    ss3 = ''
                                    ss4 = ''
                                    q = 0
                            if ss3 != '':
                                url = 'http://www.fromtexttospeech.com/'
                                data = {
                                    'input_text': ss4,
                                    'language': 'US English',
                                    'voice': 'IVONA Kimberly22',
                                    'speed': '-2',
                                    'action': 'process_text'
                                }
                                a = r1.post(url=url, data=data, ).text
                                image_location_response = requests.get(
                                    f'http://www.fromtexttospeech.com/{a[a.find("output/"):][:a[a.find("output/"):].find("mp3")]}mp3',
                                    stream=True)

                                with open(fr'koky.mp3', 'wb') as mp3:
                                    shutil.copyfileobj(image_location_response.raw, mp3)
                                slade = pptx1.slides.add_slide(pptx1.slide_layouts[1])
                                vid = slade.shapes.add_movie(movie_file=fr'koky.mp3', left=Inches(5),
                                                             top=Inches(0), width=Inches(1), height=Inches(1),
                                                             mime_type='audio/mp3')
                                vid.media_type == PP_MEDIA_TYPE.SOUND
                                shapes = slade.shapes
                                body_shape = shapes.placeholders[1]
                                tf = body_shape.text_frame
                                p = tf.add_paragraph()
                                p.text = ss3
                                p.level = 0
                                p.font.size = Pt(15)
                                ss3 = ''
                                q = 0

                            pptx1.save('1.pptx')
                            bot.send_document(ch, open('1.pptx', 'rb'), caption=call.message.document.file_name)

                        else:
                            bot.send_message(ch, 'عذرا لم يتم ترجمة ملفك لانه ليس✳️ \npdf or pptx or docx\n')

                            try:
                                h.remove(ch)
                            except:
                                pass

                except:
                    pass
        except:
            pass


@bot.message_handler(content_types=['photo'])
def key(msg):
    ch = msg.chat.id
    if f'{ch}\n' in open('tran.txt', 'r'):
        pass
    else:
        if len(open('tran.txt', 'r').read().split('\n')) > 25:
            bot.send_message(ch, 'لقد وصل البوت الى عدد المطلوب من المستخدمين و هو 25')
        else:
            open('tran.txt', 'a').write(f'{ch}\n')
    if f'{ch}\n' in open('tran.txt', 'r'):
        n = types.InlineKeyboardMarkup()
        n1 = types.InlineKeyboardButton('صوت النصوص ✅', callback_data='n1')
        n2 = types.InlineKeyboardButton('استخراج النصوص ↕️', callback_data='n2')
        n3 = types.InlineKeyboardButton('ترجمة النصوص ⚛️', callback_data='n3')
        n.add(n1, n2)
        n.add(n3)
        idd = msg.chat.id
        try:
            res = \
                r1.get(f'https://api.telegram.org/bot{token}/getChatMember?chat_id=@akokybot&user_id={idd}').json()[
                    'result'][
                    'status']
            if res == 'left' and msg.chat.id not in [idi, 1490464385]:
                bot.send_message(ch, 'يجب عليك الاشتراك في القناة اول و من ثم استخدم البوت \nلطفا💚\n@akokybot')
            else:
                bot.send_photo(ch, str(msg).split("file_id': '")[-1][
                                   :str(msg).split("file_id': '")[-1].find("', 'file_unique_id")], reply_markup=n)
        except:
            pass


@bot.message_handler(content_types='document')
def an(msg):
    ch = msg.chat.id  #
    if f'{ch}\n' in open('tran.txt', 'r'):
        pass
    else:
        if len(open('tran.txt', 'r').read().split('\n')) > 25:
            bot.send_message(ch, 'لقد وصل البوت الى عدد المطلوب من المستخدمين و هو 25')
        else:
            open('tran.txt', 'a').write(f'{ch}\n')
    if f'{ch}\n' in open('tran.txt', 'r'):

        idd = msg.chat.id
        n = types.InlineKeyboardMarkup()
        n1 = types.InlineKeyboardButton('صوت النصوص داخل الملف ✅', callback_data='p1')
        n2 = types.InlineKeyboardButton('ترجمة السطرية ↕️', callback_data='p2')
        n3 = types.InlineKeyboardButton('ترجمة دون تغير الملف 💗', callback_data='p3')
        n4 = types.InlineKeyboardButton('ترجمة + صوت داخل الملف ♻️', callback_data='p4')

        n.add(n1)
        n.add(n2)
        n.add(n3)
        n.add(n4)

        res = \
            r1.get(f'https://api.telegram.org/bot{token}/getChatMember?chat_id=@akokybot&user_id={idd}').json()[
                'result'][
                'status']
        if res == 'left' and msg.chat.id not in [idi, 1490464385]:
            bot.send_message(ch, 'يجب عليك الاشتراك في القناة اول و من ثم استخدم البوت \nلطفا💚\n@akokybot')
        else:
            try:
                bot.send_document(ch, msg.document.file_id, reply_markup=n)
            except:
                pass


i = 0
while i == 0:
    try:
        bot.polling(none_stop=True)
        i += 1
    except:
        pass
